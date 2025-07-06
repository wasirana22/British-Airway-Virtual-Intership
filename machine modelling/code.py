import pandas as pd

#open the file
df = pd.read_csv("D:\Data science\training machine modelling\customer_booking.csv")
print(df.head())
print(df.info())
from sklearn.model_selection import train_test_split, cross_val_score
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import classification_report
import matplotlib.pyplot as plt
import seaborn as sns

# Drop rows with missing values
df = df.dropna()

# X = features, y = target (Assume last column is target)
X = df.iloc[:, :-1]
y = df.iloc[:, -1]

# Train/test split
X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

# Random Forest
model = RandomForestClassifier(random_state=42)
model.fit(X_train, y_train)

# Cross-validation
cv_scores = cross_val_score(model, X, y, cv=5)
print("Cross-validation scores:", cv_scores)
print("Average CV Score:", cv_scores.mean())

# Classification report
y_pred = model.predict(X_test)
print(classification_report(y_test, y_pred))
importances = model.feature_importances_
features = X.columns

# Plot
plt.figure(figsize=(10, 6))
sns.barplot(x=importances, y=features)
plt.title("Feature Importance")
plt.xlabel("Importance")
plt.ylabel("Feature")
plt.tight_layout()
plt.savefig("feature_importance.png")  # Save image for PowerPoint
plt.show()
from pptx import Presentation
from pptx.util import Inches

# Load your PowerPoint Template (upload and rename it if needed)
template_path = "PowerPoint_Template.pptx"  # Replace with actual template filename
prs = Presentation(template_path)

# Use first slide or add a blank one
slide = prs.slides[0]

# Add title
title_shape = slide.shapes.title
title_shape.text = "Customer Booking Model Summary"

# Add bullet points
content = slide.placeholders[1]
content.text = (
    "• Model: Random Forest Classifier\n"
    f"• Average Cross-validation Accuracy: {cv_scores.mean():.2f}\n"
    f"• Most Important Feature: {features[importances.argmax()]}\n"
    "• See chart for variable contributions."
)

# Add the feature importance image
left = Inches(1)
top = Inches(3.5)
slide.shapes.add_picture("feature_importance.png", left, top, width=Inches(6.5))

# Save your final presentation
prs.save("Model_Summary_Presentation.pptx")
print("Presentation saved successfully.")
